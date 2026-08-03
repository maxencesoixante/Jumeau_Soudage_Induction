"""Applique les commentaires PowerPoint sur le deck live (édition chirurgicale).

NE PAS confondre avec build_deck.py (qui reconstruit le deck et n'est PAS
idempotent). Ce script, lui, modifie le deck existant :
  1. sauvegarde d'abord une copie datée ;
  2. remplace 5 figures régénérées (blob des parts image) ;
  3. re-ajuste les 2 schémas (nouveau format paysage côte à côte) ;
  4. renomme CFC -> MFC dans le texte des slides ;
  5. ajoute un bullet « spot x=60 » sur la slide dissipation ;
  6. insère une nouvelle slide « profil en M + loi R=k·I²−L » après la slide 9.
"""
import copy
import shutil
import datetime as dt
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

R = Path("/Users/maxencedubois/PycharmProjects/Jumeau_Soudage_Induction")
DECK = R / "docs" / "Point d'avancement hebdomadaire — LIPeC  ÉTS.pptx"
FIG = R / "docs" / "figures"
BODY = RGBColor(0x3A, 0x3A, 0x3A)
GREY = RGBColor(0x9A, 0x9A, 0x9A)
ACCENT = RGBColor(0xC1, 0x27, 0x2D)
SHARED_TITLE = "Jumeau numérique — modélisation & essais de validation"

# --- 1. sauvegarde -----------------------------------------------------
bak = DECK.with_name(DECK.stem + f".backup-{dt.date.today().isoformat()}.pptx")
if not bak.exists():
    shutil.copy2(DECK, bak)
    print("backup:", bak.name)

prs = Presentation(str(DECK))
slides = prs.slides

# --- 2/3. remplacement + re-fit des figures ----------------------------
REPL = {
    "image12.png": FIG / "schema_montage_exp7.png",
    "image15.png": FIG / "fig2_mesure_modele.png",
    "image17.png": FIG / "schema_montage_exp9.png",
    "image18.png": FIG / "fig_dissipation_monospot.png",
    "image19.png": FIG / "fig_dissipation_semistatique.png",
}
REFIT = {"image12.png", "image17.png"}   # aspect changé (paysage) → re-fit

n_repl = 0
for s in slides:
    for sh in s.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        rId = sh._element.blipFill.blip.rEmbed
        part = sh.part.related_part(rId)
        key = part.partname.split("/")[-1]
        if key not in REPL:
            continue
        png = REPL[key]
        with open(png, "rb") as f:
            part._blob = f.read()
        n_repl += 1
        if key in REFIT:
            iw, ih = Image.open(png).size
            ar = iw / ih
            bx, by, bw, bh = sh.left, sh.top, sh.width, sh.height
            if ar > bw / bh:
                nw, nh = bw, int(round(bw / ar))
            else:
                nh, nw = bh, int(round(bh * ar))
            sh.left = bx + (bw - nw) // 2
            sh.top = by + (bh - nh) // 2
            sh.width = nw
            sh.height = nh
print("images remplacées:", n_repl)

# --- 4. CFC -> MFC dans le texte des slides ----------------------------
n_txt = 0
for s in slides:
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for run in p.runs:
                if "CFC" in run.text:
                    run.text = run.text.replace("CFC", "MFC")
                    n_txt += 1
print("runs CFC→MFC:", n_txt)


# --- helpers texte -----------------------------------------------------
def find(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_single(shape, text):
    p = shape.text_frame.paragraphs[0]
    p.runs[0].text = text
    for extra in p.runs[1:]:
        extra._r.getparent().remove(extra._r)


def add_bullet(shape, text, size=20.25, color=BODY, bullet="•  "):
    tf = shape.text_frame
    tf.word_wrap = True
    para = tf.add_paragraph()
    para.space_after = Pt(10)
    run = para.add_run()
    run.text = f"{bullet}{text}"
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = False
    run.font.color.rgb = color


# --- 5. bullet « spot x=60 » sur la slide 11 (dissipation) -------------
s11 = slides[10]
tb = find(s11, "TextBox 17")
if tb is not None and "x = 60" not in tb.text_frame.text:
    add_bullet(tb, "Le spot est centré à x = 60 mm (centre du coil) → c'est là que "
                   "la température de pic est maximale (TC3).", size=18)
    print("bullet x=60 ajouté")


# --- 6. nouvelle slide « profil en M + loi courant » -------------------
def duplicate(src):
    new = slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    rid_map = {}
    for rid, rel in src.part.rels.items():
        rt = rel.reltype
        if rt.endswith("slideLayout") or rt.endswith("notesSlide"):
            continue
        new_rid = (new.part.relate_to(rel._target, rt, is_external=True)
                   if rel.is_external else new.part.relate_to(rel._target, rt))
        rid_map[rid] = new_rid
    for el in new.shapes._spTree.iter():
        for attr in (qn("r:embed"), qn("r:link"), qn("r:id")):
            v = el.get(attr)
            if v in rid_map:
                el.set(attr, rid_map[v])
    return new


def textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def heading(slide, x, y, w, text, color=ACCENT, size=21):
    tb = textbox(slide, x, y, w, 0.7)
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.name, r.font.size, r.font.bold = "Arial", Pt(size), True
    r.font.color.rgb = color
    return tb


def bullets_box(slide, x, y, w, h, lines, size=16.5):
    tb = textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(9)
        run = para.add_run()
        run.text = f"•  {line}"
        run.font.name, run.font.size = "Arial", Pt(size)
        run.font.color.rgb = BODY
    return tb


# cloner la slide 9 (modélisation, en-tête partagé) et la vider
tmpl = slides[8]
new = duplicate(tmpl)
for sh in list(new.shapes):
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        sh._element.getparent().remove(sh._element)
    elif sh.name == "TextBox 19":            # ancien bloc de bullets
        sh._element.getparent().remove(sh._element)

# en-tête
if find(new, "TextBox 12"):
    set_single(find(new, "TextBox 12"), "MODÉLISATION SUR PYTHON")
if find(new, "TextBox 13"):
    set_single(find(new, "TextBox 13"), SHARED_TITLE)
if find(new, "TextBox 14"):
    set_single(find(new, "TextBox 14"), "Comment le modèle génère le profil en M & la loi en courant")

# colonne gauche : physique du M
heading(new, 1.2, 2.95, 9.0, "Comment le modèle génère le profil en « M »")
bullets_box(new, 1.2, 3.75, 8.9, 7.0, [
    "Le champ EM (bobine + MFC, Biot–Savart) induit des courants de Foucault "
    "dans le pli twill conducteur.",
    "Plaque mince devant l'épaisseur de peau (δ ≈ 6 mm > 3,36 mm) → courants "
    "plans, champ Bz quasi uniforme (Lin 1993 ; Grouve 2020).",
    "Le solveur résout la fonction de courant ψ avec ψ = 0 au bord : le courant "
    "ne traverse pas le chant, les boucles se referment dans la plaque.",
    "La densité J = ∇×(ψ ẑ) s'accumule vers les bords → puissance Joule ∝ ρ·J² : "
    "deux lobes chauds aux chants, centre plus froid = le « M ».",
    "Renforcé géométriquement : bobine + MFC (55 mm) débordent la largeur du "
    "coupon (40 mm).",
])

# colonne droite : origine de la loi en courant
heading(new, 10.7, 2.95, 8.6, "D'où vient la loi  R = k·I² − L ?")
bullets_box(new, 10.7, 3.75, 8.5, 7.0, [
    "Source Joule ∝ Bz² et Bz ∝ I (Biot–Savart) → puissance déposée ∝ I².",
    "Taux de chauffe initial dT/dt ∝ puissance → terme moteur k·I².",
    "Pertes (convection + conduction vers le puits céramique/MFC) quasi "
    "constantes → terme −L (~3,5 °C/s).",
    "Ajustement mesuré : R² = 0,999 sur 5 courants ; fréquence machine "
    "constante 388 ± 2 kHz (pas de couplage fréquence–courant).",
    "Références : Lin 1993 ; Grouve 2020 ; O'Shaughnessey 2014 ; "
    "Fluxtrol Ferrotron 559H.",
])

# --- placer la nouvelle slide juste après la slide 9 -------------------
sldIdLst = slides._sldIdLst
ids = list(sldIdLst)
new_id = ids[-1]                 # la slide ajoutée est en dernier
sldIdLst.remove(new_id)
sldIdLst.insert(9, new_id)       # index 9 = juste après la 9e slide (0-based 8)

prs.save(str(DECK))
print("saved:", DECK.name, "| n slides:", len(list(sldIdLst)))
